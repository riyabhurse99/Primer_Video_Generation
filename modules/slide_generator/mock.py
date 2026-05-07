from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from models.schemas import VideoScript
from modules.slide_generator.base import BaseSlideGenerator
from utils.logger import get_logger

logger = get_logger(__name__)

# Scaler brand colours
SCALER_DARK = RGBColor(0x1A, 0x1A, 0x2E)
SCALER_BLUE = RGBColor(0x16, 0x21, 0x3E)
SCALER_ACCENT = RGBColor(0x0F, 0x3D, 0x87)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


class MockSlideGenerator(BaseSlideGenerator):
    """
    Generates a real PPTX using python-pptx.
    No API key needed — fully functional.
    Replace with OpenMAICSlideGenerator once OpenMAIC is set up.
    """

    def _set_slide_background(self, slide, color: RGBColor):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_title_slide(self, prs: Presentation, topic: str, depth: str):
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, SCALER_DARK)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = topic
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.8))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"Level: {depth.capitalize()}  |  Scaler Primer"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)

    def _add_content_slide(self, prs: Presentation, title: str, content: list[str]):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, SCALER_BLUE)

        # Title bar
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.9))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.15),
            Inches(9.0), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = SCALER_ACCENT
        line.line.fill.background()

        # Content bullets
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.0))
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        for i, point in enumerate(content):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = f"•  {point}"
            p.font.size = Pt(22)
            p.font.color.rgb = WHITE
            p.space_after = Pt(12)

    def generate(self, video_script: VideoScript, output_path: str) -> str:
        logger.info(f"[MOCK] Generating PPTX for: {video_script.topic}")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        self._add_title_slide(prs, video_script.topic, video_script.depth)

        for slide in video_script.slides:
            self._add_content_slide(prs, slide.title, slide.content)

        prs.save(output_path)
        logger.info(f"[MOCK] PPTX saved: {output_path}")
        return output_path
