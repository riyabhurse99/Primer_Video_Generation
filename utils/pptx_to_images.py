import os
import shutil
import subprocess
import platform
from utils.logger import get_logger

logger = get_logger(__name__)


def _find_libreoffice_cmd() -> str:
    """Find the correct LibreOffice command for the current OS."""
    system = platform.system()
    if system == "Darwin":  # macOS
        mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.exists(mac_path):
            return mac_path
        if shutil.which("soffice"):
            return "soffice"
    elif system == "Linux":
        if shutil.which("libreoffice"):
            return "libreoffice"
        if shutil.which("soffice"):
            return "soffice"
    elif system == "Windows":
        if shutil.which("soffice"):
            return "soffice"
    return "libreoffice"


def _has_tool(cmd: str) -> bool:
    return shutil.which(cmd) is not None


def _convert_with_libreoffice(pptx_path: str, output_dir: str) -> list[str]:
    """Full quality: PPTX → PDF (LibreOffice) → PNG (pdftoppm)."""
    lo_cmd = _find_libreoffice_cmd()
    logger.info(f"Converting PPTX to PDF: {pptx_path} (using: {lo_cmd})")

    result = subprocess.run(
        [lo_cmd, "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
    pdf_path = os.path.join(output_dir, pdf_name)

    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF not created at expected path: {pdf_path}")

    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    output_prefix = os.path.join(output_dir, base_name)

    logger.info(f"Converting PDF to PNG images: {pdf_path}")
    result = subprocess.run(
        ["pdftoppm", "-png", "-r", "150", pdf_path, output_prefix],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        raise RuntimeError(f"pdftoppm conversion failed: {result.stderr}")

    images = sorted([
        os.path.join(output_dir, f)
        for f in os.listdir(output_dir)
        if f.startswith(base_name) and f.endswith(".png")
    ])

    if not images:
        raise FileNotFoundError("No PNG images were generated from the PPTX")

    logger.info(f"Generated {len(images)} slide images (LibreOffice method)")
    return images


def _convert_with_pillow_fallback(pptx_path: str, output_dir: str) -> list[str]:
    """
    Fallback: renders placeholder slide images using Pillow.
    Used when LibreOffice and poppler are not installed.
    Reads slide content from the PPTX and creates text-rendered PNGs.
    """
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont

    logger.info(f"Using Pillow fallback to render slides from: {pptx_path}")
    prs = Presentation(pptx_path)
    images = []

    # Try to use a monospace font, fall back to default
    try:
        font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 48)
        font_body = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
    except (IOError, OSError):
        try:
            font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
            font_body = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
        except (IOError, OSError):
            font_title = ImageFont.load_default()
            font_body = ImageFont.load_default()

    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (1920, 1080), color=(22, 33, 62))  # Scaler dark blue
        draw = ImageDraw.Draw(img)

        # Extract text from the slide
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

        # Draw title
        y_pos = 80
        if texts:
            draw.text((80, y_pos), texts[0], font=font_title, fill=(255, 255, 255))
            y_pos += 100

            # Draw line separator
            draw.line([(80, y_pos), (1840, y_pos)], fill=(15, 61, 135), width=3)
            y_pos += 40

            # Draw remaining content
            for text in texts[1:]:
                draw.text((100, y_pos), text, font=font_body, fill=(220, 220, 220))
                y_pos += 55
                if y_pos > 950:
                    break

        img_path = os.path.join(output_dir, f"slide_{i:03d}.png")
        img.save(img_path)
        images.append(img_path)

    logger.info(f"Generated {len(images)} slide images (Pillow fallback)")
    return images


def pptx_to_images(pptx_path: str, output_dir: str) -> list[str]:
    """
    Converts each slide in a PPTX to a PNG image.
    Tries LibreOffice + poppler first (best quality).
    Falls back to Pillow-based rendering if tools not installed.
    Returns list of image paths in slide order.

    Special case: if a .png_list file exists alongside the PPTX (written by
    GrootSlideGenerator), those pre-rendered PNG paths are returned directly,
    skipping all conversion.
    """
    import json

    # Groot proxy: pre-rendered PNGs already exist
    png_list_path = pptx_path.replace(".pptx", ".png_list")
    if os.path.exists(png_list_path):
        with open(png_list_path) as f:
            paths = json.load(f)
        logger.info(f"Using pre-rendered PNG list ({len(paths)} slides): {png_list_path}")
        return paths

    os.makedirs(output_dir, exist_ok=True)

    # Try the high-quality path first
    lo_cmd = _find_libreoffice_cmd()
    has_lo = shutil.which(lo_cmd) is not None or os.path.exists(lo_cmd)
    has_pdftoppm = _has_tool("pdftoppm")

    if has_lo and has_pdftoppm:
        try:
            return _convert_with_libreoffice(pptx_path, output_dir)
        except Exception as e:
            logger.warning(f"LibreOffice method failed: {e}. Falling back to Pillow.")

    # Fallback
    return _convert_with_pillow_fallback(pptx_path, output_dir)
